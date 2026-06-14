from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import build_html_dashboard as dash


OUT = Path("outputs/negative_feedback_af_summary.pptx")

WIDE_W = 13.333
WIDE_H = 7.5

INK = RGBColor(20, 54, 66)
TEAL = RGBColor(31, 122, 122)
RED = RGBColor(194, 65, 45)
ORANGE = RGBColor(225, 129, 44)
YELLOW = RGBColor(216, 166, 37)
BLUE = RGBColor(49, 93, 133)
MUTED = RGBColor(100, 116, 139)
PAPER = RGBColor(246, 248, 250)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(216, 224, 231)


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Tahoma"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_kicker(slide, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.35), Inches(0.12), Inches(0.12))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    add_text(slide, 0.76, 0.26, 4.4, 0.28, text.upper(), 8.5, True, TEAL)


def add_title(slide, kicker, title, subtitle=""):
    add_kicker(slide, kicker)
    add_text(slide, 0.55, 0.66, 11.7, 0.78, title, 30, True, INK)
    if subtitle:
        add_text(slide, 0.58, 1.38, 11.4, 0.42, subtitle, 13, False, MUTED)


def add_footer(slide, idx):
    add_text(slide, 0.55, 7.1, 7.2, 0.2, "Negative Feedback: ไม่มีซิม True ดู AF ไม่ได้", 7.5, False, MUTED)
    add_text(slide, 12.25, 7.08, 0.55, 0.22, f"{idx:02d}", 8, True, MUTED, PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, label, value, accent=TEAL, note=""):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    add_text(slide, x + 0.15, y + 0.14, w - 0.3, 0.24, label, 9.5, True, MUTED)
    add_text(slide, x + 0.15, y + 0.48, w - 0.3, 0.42, str(value), 23, True, accent)
    if note:
        add_text(slide, x + 0.15, y + 0.96, w - 0.3, 0.34, note, 8.8, False, MUTED)


def add_bar(slide, x, y, w, label, value, max_value, color=RED, suffix=""):
    add_text(slide, x, y - 0.02, 2.0, 0.22, label, 10, True, INK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 2.15), Inches(y), Inches(w), Inches(0.16))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(235, 240, 244)
    bg.line.fill.background()
    bar_w = 0 if max_value == 0 else w * value / max_value
    fg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 2.15), Inches(y), Inches(bar_w), Inches(0.16))
    fg.fill.solid()
    fg.fill.fore_color.rgb = color
    fg.line.fill.background()
    add_text(slide, x + 2.25 + w, y - 0.06, 1.2, 0.26, f"{value}{suffix}", 9.5, True, INK, PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, h, headers, rows, widths=None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
        for p in cell.text_frame.paragraphs:
            p.runs[0].font.name = "Tahoma"
            p.runs[0].font.size = Pt(8.5)
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = WHITE
    for r, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.runs[0].font.name = "Tahoma"
                p.runs[0].font.size = Pt(8)
                p.runs[0].font.color.rgb = INK
    return table_shape


def clip(text: str, n: int = 210) -> str:
    text = " ".join(str(text).split())
    return text if len(text) <= n else text[: n - 1] + "…"


def build_deck():
    records = pd.concat([dash.load_original_records(), dash.load_new_records()], ignore_index=True)
    summary = dash.build_summary(records)
    neg = records[records["is_negative_final"] == "Yes"].copy()
    platform_rows = summary["platforms"]
    leaders = summary["leaders"][:6]
    keywords = summary["keywords"][:8]
    fb = next(row for row in platform_rows if row["platform"] == "Facebook")
    xrow = next(row for row in platform_rows if row["platform"] == "X")

    prs = Presentation()
    prs.slide_width = Inches(WIDE_W)
    prs.slide_height = Inches(WIDE_H)

    blank = prs.slide_layouts[6]

    # Slide 1
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    add_kicker(s, "Executive readout")
    add_text(s, 0.55, 0.85, 11.6, 1.05, "Negative feedback กระจุกอยู่ที่ปัญหาดู AF ไม่ได้ ไม่ใช่กระแสทั่วไป", 32, True, INK)
    add_text(s, 0.58, 1.9, 10.9, 0.45, "หลังตัดโพสต์โปรโมต/กิจกรรมออก เหลือ feedback ที่เกี่ยวกับดูไม่ได้, ซิม True, แพ็กเกจ และเงื่อนไขรับชมโดยตรง", 14, False, MUTED)
    add_card(s, 0.6, 3.0, 2.2, 1.45, "Relevant negative", summary["negative"], RED, f"{summary['negative_share']:.1%} of mentions")
    add_card(s, 3.1, 3.0, 2.2, 1.45, "Facebook", fb["negative"], RED, "ช่องทางหลักของปัญหา")
    add_card(s, 5.6, 3.0, 2.2, 1.45, "X", xrow["negative"], ORANGE, "ข้อวิจารณ์คมและแชร์ง่าย")
    add_card(s, 8.1, 3.0, 2.2, 1.45, "Access issue", summary["access"], TEAL, "ดูไม่ได้ / ช่องทางรับชม")
    add_card(s, 10.6, 3.0, 2.2, 1.45, "True/package", summary["true_package"], BLUE, "ซิม / แพ็กเกจ / ลูกค้า True")
    add_text(s, 0.7, 5.05, 11.5, 0.7, "สัญญาณหลัก: คนไม่ได้โต้แย้งตัวรายการ แต่โต้แย้ง friction ของการรับชมคอนเสิร์ตและเงื่อนไขที่ทำให้คนที่อยากดูเข้าไม่ถึง", 20, True, INK)
    add_footer(s, 1)

    # Slide 2
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_title(s, "Channel pressure", "Facebook เป็นศูนย์กลาง negative feedback จากปัญหาดูไม่ได้", "ตัวเลขนี้นับเฉพาะ negative ที่เกี่ยวข้องกับการรับชม/ซิม/แพ็กเกจ ไม่รวมโพสต์โปรโมตทั่วไป")
    max_neg = max(row["negative"] for row in platform_rows)
    y = 2.15
    for row in platform_rows:
        add_bar(s, 0.75, y, 6.5, row["platform"], row["negative"], max_neg, RED, f" ({row['negative_share']:.0%})")
        y += 0.48
    add_text(s, 9.7, 2.0, 2.2, 0.3, "Readout", 14, True, TEAL)
    add_text(s, 8.55, 2.42, 3.55, 1.0, f"Facebook คิดเป็น {fb['negative']} mentions จาก relevant negative ทั้งหมด {summary['negative']} mentions", 18, True, INK)
    add_text(s, 8.55, 3.62, 3.55, 1.25, "X มี volume น้อยกว่า แต่ข้อความมักเป็นคำวิจารณ์ตรง เช่น สมัครแล้วดูไม่ได้, ต้องใช้ True/dtac, เงื่อนไขไม่ชัด", 12.5, False, MUTED)
    add_footer(s, 2)

    # Slide 3
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    add_title(s, "Issue anatomy", "Pain point หลักคือ access friction มากกว่า sentiment ต่อตัวรายการ", "")
    add_card(s, 0.7, 2.0, 2.7, 1.55, "ช่องทางรับชม / ดูไม่ได้", summary["access"], RED, "core pain point")
    add_card(s, 3.75, 2.0, 2.7, 1.55, "ซิม True / แพ็กเกจ", summary["true_package"], ORANGE, "condition friction")
    add_card(s, 6.8, 2.0, 2.7, 1.55, "ดูไม่ได้ / ต้องสมัคร", summary["unable"], TEAL, "explicit access failure")
    add_card(s, 9.85, 2.0, 2.7, 1.55, "ถ้อยคำต่อว่า", summary["complaint"], BLUE, "strong wording")
    add_text(s, 0.85, 4.25, 5.6, 0.8, "ตีความ: คนมี intent จะดูและติดตามอยู่แล้ว แต่เจอ friction ระหว่างจุดที่ควร convert เป็น engagement/โหวต", 20, True, INK)
    add_text(s, 7.1, 4.25, 4.75, 1.3, "Implication: ช่องทางรับชมควรสื่อสารแบบ zero ambiguity ก่อนคอนเสิร์ต และแยก free access / paid package / telco condition ให้ชัด", 13, False, MUTED)
    add_footer(s, 3)

    # Slide 4
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_title(s, "Language signal", "คำลบวนอยู่กับ ‘ดูไม่ได้’ และ ‘เงื่อนไขที่ทำให้เข้าไม่ถึง’", "")
    max_kw = max(k["count"] for k in keywords) if keywords else 1
    y = 2.0
    for kw in keywords:
        add_bar(s, 0.8, y, 7.5, kw["keyword"], kw["count"], max_kw, TEAL)
        y += 0.5
    add_text(s, 9.45, 2.0, 2.9, 0.35, "What it means", 14, True, TEAL)
    add_text(s, 8.8, 2.45, 3.45, 1.75, "Keyword ไม่ได้บอกว่าแฟนไม่ชอบรายการ แต่บอกว่าเขาสะดุดกับวิธีดู: ดูไม่ได้, ดูยาก, ต้องใช้, ไม่ดู, เสียความรู้สึก", 16, True, INK)
    add_text(s, 8.8, 4.55, 3.45, 0.9, "นี่เป็น operational issue ที่แก้ได้ด้วย product access + communication มากกว่าการแก้ content", 12.5, False, MUTED)
    add_footer(s, 4)

    # Slide 5
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    add_title(s, "Amplification accounts", "กระแสลบถูกขยายโดย account/เพจที่มี engagement หรือถ้อยคำแรง", "ตารางนี้ตัด owned brand account และ UNKNOWN ออก เพื่อดู account ที่พากระแสชัดขึ้น")
    rows = []
    for row in leaders:
        rows.append([row["platform"], row["account"], row["negative_posts"], row["total_negative_engagement"], int(row["impact_score"])])
    add_table(s, 0.65, 1.92, 12.0, 3.25, ["Platform", "Account", "Neg.", "Eng.", "Impact"], rows, [1.4, 5.5, 1.0, 1.4, 1.2])
    add_text(s, 0.75, 5.55, 11.7, 0.45, "Key read: เพจ ‘อวยไส้แตกแหกไส้ฉีก’ เป็นตัวอย่าง account ที่ทำให้ประเด็นดูไม่ได้ถูกเล่าเป็น narrative ยาวและแชร์ต่อได้", 16, True, INK)
    add_footer(s, 5)

    # Slide 6
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_title(s, "Voice of customer", "ตัวอย่าง feedback สะท้อนความเสียหายที่เกิดตอนคนพร้อมจะดู", "")
    examples = neg.sort_values(["negative_score", "engagement"], ascending=False).head(4)
    x_positions = [0.65, 6.75, 0.65, 6.75]
    y_positions = [1.9, 1.9, 4.35, 4.35]
    for i, (_, row) in enumerate(examples.iterrows()):
        x, y = x_positions[i], y_positions[i]
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.65), Inches(1.65))
        rect.fill.solid()
        rect.fill.fore_color.rgb = PAPER
        rect.line.color.rgb = LINE
        add_text(s, x + 0.18, y + 0.13, 5.2, 0.24, f"{row['platform']} · {row['account']} · score {row['negative_score']}", 8.5, True, RED)
        add_text(s, x + 0.18, y + 0.45, 5.25, 1.0, clip(row["message"], 190), 10.3, False, INK)
    add_footer(s, 6)

    # Slide 7
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill_fore_color = PAPER
    s.background.fill.fore_color.rgb = PAPER
    add_title(s, "Action takeaway", "ลด negative ได้ด้วยการลด ambiguity ก่อนถึงจุดรับชม", "")
    actions = [
        ("1", "ประกาศ matrix การรับชม", "ฟรี/เสียเงิน/ต้องเป็นลูกค้า/ต้องใช้ซิม ใส่เป็นตารางเดียวก่อนคอน"),
        ("2", "ทำ fallback path", "ถ้าเข้า app/ช่องหลักไม่ได้ ต้องมี link หรือ step สำรองที่ชัด"),
        ("3", "สื่อสารก่อนคอนเสิร์ต", "โพสต์ reminder พร้อมตัวอย่างหน้าจอ และ pin วิธีดูในทุก platform"),
        ("4", "จับ complaint realtime", "ระหว่างคอนให้มีทีมตอบปัญหา ‘ดูไม่ได้/กดไม่ได้/สมัครแล้วไม่ได้’ ทันที"),
    ]
    for i, (num, head, body) in enumerate(actions):
        y = 1.8 + i * 1.08
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.82), Inches(y), Inches(0.42), Inches(0.42))
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        add_text(s, 0.82, y + 0.07, 0.42, 0.18, num, 9, True, WHITE, PP_ALIGN.CENTER)
        add_text(s, 1.45, y - 0.02, 3.6, 0.28, head, 15, True, INK)
        add_text(s, 5.0, y - 0.02, 6.8, 0.32, body, 12.5, False, MUTED)
    add_text(s, 0.8, 6.45, 11.2, 0.34, "Bottom line: ปัญหานี้แก้ได้เร็วถ้าทำให้เส้นทางการดู ‘เข้าใจง่ายที่สุด’ ก่อนเกิดคอนเสิร์ตครั้งถัดไป", 15, True, RED)
    add_footer(s, 7)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    output = build_deck()
    print(output)
